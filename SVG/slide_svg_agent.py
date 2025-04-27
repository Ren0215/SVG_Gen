"""
slide_svg_agent.py – グローバル文脈確実修正版
=============================
- グローバル文脈の確実表示を保証
- Pydanticモデルのエイリアス問題を回避
- LangGraphとLLM機能を維持
"""
from __future__ import annotations
import argparse, json, os, re, time
from pathlib import Path
from typing import List, Dict, Any, Literal, Optional
from dotenv import load_dotenv

import jinja2
from pydantic import BaseModel, Field
from langgraph.graph import StateGraph, END
from openai import OpenAI

# 環境変数読み込み
load_dotenv()

# ──────────── Paths & env ────────────
ROOT = Path(__file__).parent
TPL_DIR = ROOT / "templates"
ENV = jinja2.Environment(loader=jinja2.FileSystemLoader(str(TPL_DIR)))
ENV.globals.update(len=len, enumerate=enumerate, zip=zip)
ENV.filters["split"] = lambda s, sep="\n": s.split(sep)

# LLM初期化
try:
    OAI = OpenAI()
    OPENAI_AVAILABLE = True
except Exception as e:
    print(f"OpenAI初期化エラー: {e}")
    OPENAI_AVAILABLE = False

# ──────────── State schema ────────────
class SlidesState(BaseModel):
    raw_json: Dict
    company: str | None = None
    # カテゴリーごとのデータ
    strategy: List[str] = []
    product: List[str] = []
    top: List[str] = []
    industry: List[str] = []
    society: List[str] = []
    # グローバル文脈用フィールド（key=valueの形でバックアップも保持）
    global_items: List[str] = []  # 名前を変更して問題を回避
    # その他の設定
    template_id: str = "context_map_dark"
    palette: str = "context_map_dark"
    svg: str | None = None
    status: str | None = None

# ──────────── パース用定数 ────────────
CAT_MAP = {
    "企業戦略文脈": "strategy",
    "プロダクト文脈": "product",
    "TOP文脈": "top", 
    "業界トレンド文脈": "industry",
    "社会トレンド文脈": "society",
    "グローバル文脈": "global_items",  # 名前を変更
}
URL_RE = re.compile(r"→?https?://\S+")
HDR_RE = re.compile(r"^[\s　\r\n]*[（\(]?[A-FＡ-Ｆ][\)）]\s*(.+?文脈)")
BULLET_RE = re.compile(r"^[\-\*‧•●◦・]\s*")

# ──────────── LLM機能 ────────────────
def summarize_with_openai(text: str, category: str) -> str:
    """OpenAIを使って文言を要約する"""
    if not OPENAI_AVAILABLE or os.environ.get("SKIP_LLM", "").lower() == "true":
        # LLM利用不可または無効化されている場合は簡易要約を返す
        return text[:20] + ("…" if len(text) > 20 else "")
    
    # カテゴリに応じた文脈を設定
    context_map = {
        "strategy": "企業戦略に関する内容",
        "product": "製品・サービスに関する内容",
        "top": "経営トップに関する内容",
        "industry": "業界トレンドに関する内容",
        "society": "社会トレンドに関する内容",
        "global_items": "グローバル動向に関する内容"
    }
    context = context_map.get(category, "一般的な内容")
    
    # 最大文字数を設定
    max_length = 20  # 楕円ノードに収まる文字数
    
    # OpenAI GPT-4による要約（リトライ処理付き）
    for attempt in range(3):  # 最大3回リトライ
        try:
            prompt = f"""
            以下のテキストを、最大{max_length}文字以内の簡潔な表現に要約してください。
            元のテキスト: 「{text}」
            カテゴリ: {context}
            
            条件:
            - {max_length}文字以内に収める
            - 元の意味を保持する
            - 専門用語や固有名詞は維持する
            - 文末の「」や句読点は省略可能
            - キーワードを含める
            
            要約のみを出力してください。
            """
            
            response = OAI.chat.completions.create(
                model="o4-mini",
                messages=[{"role": "user", "content": prompt}]
            )
            
            summary = response.choices[0].message.content.strip()
            print(f"OpenAI要約: '{text}' → '{summary}'")
            return summary[:max_length]
            
        except Exception as e:
            print(f"OpenAI処理エラー (試行 {attempt+1}/3): {e}")
            if attempt < 2:  # 3回目の失敗では即座に戻る
                time.sleep(2)  # 少し待ってからリトライ
    
    # すべてのリトライが失敗した場合は簡易要約を返す
    return text[:max_length] + ("…" if len(text) > max_length else "")

# ──────────── Node functions ────────────
def parse_context(state: SlidesState) -> SlidesState:
    """コンテキストマップのテキストをパースして6カテゴリ×5行に分解"""
    extracted = {v: [] for v in CAT_MAP.values()}
    cur_cat = None
    
    # グローバル文脈データを確実に抽出するための追跡変数
    global_items = []
    
    for raw in state.raw_json["text"].splitlines():
        # テキスト前処理
        line = BULLET_RE.sub("", URL_RE.sub("", raw)).strip(' "\u3000')
        if not line:
            continue
        
        # ヘッダー検出
        if m := HDR_RE.match(line):
            category = m.group(1)
            cur_cat = CAT_MAP.get(category)
            
            # グローバル文脈ヘッダーの検出を追跡
            if category == "グローバル文脈":
                print(f"DEBUG: グローバル文脈ヘッダー検出！")
            
            continue
            
        # 本文追加
        if cur_cat and len(extracted[cur_cat]) < 5:
            extracted[cur_cat].append(line)
            
            # グローバル文脈データを追跡
            if cur_cat == "global_items":
                global_items.append(line)
                print(f"DEBUG: グローバル文脈アイテム追加: {line}")
    
    # デバッグ出力
    for k, v in extracted.items():
        print(f"DEBUG: {k} bucket = {v}")
    
    # グローバルデータ特別処理
    print(f"DEBUG: グローバル文脈追跡データ = {global_items}")
    
    # 更新ディクショナリ作成
    update_dict = {}
    
    # 個別にフィールドを更新
    update_dict["strategy"] = extracted["strategy"]
    update_dict["product"] = extracted["product"]
    update_dict["top"] = extracted["top"]
    update_dict["industry"] = extracted["industry"]
    update_dict["society"] = extracted["society"]
    update_dict["global_items"] = extracted["global_items"]  # グローバル文脈
    update_dict["template_id"] = "context_map_dark"
    update_dict["palette"] = "context_map_dark"
    
    # 更新状態を返す
    result = state.model_copy(update=update_dict)
    print(f"DEBUG: global_items after update = {result.global_items}")
    return result

def summarize_content(state: SlidesState) -> SlidesState:
    """各カテゴリの文言をLLMを使って要約"""
    # 処理が重いので、実行時引数でスキップできるようにする
    if os.environ.get("SKIP_LLM", "").lower() == "true":
        print("LLM処理をスキップします")
        # 長いテキストは手動で短くする（LLMを使わない場合）
        update_dict = {}
        
        # カテゴリごとに処理
        update_dict["strategy"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.strategy]
        update_dict["product"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.product]
        update_dict["top"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.top]
        update_dict["industry"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.industry]
        update_dict["society"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.society]
        update_dict["global_items"] = [item[:20] + ("…" if len(item) > 20 else "") for item in state.global_items]
        
        return state.model_copy(update=update_dict)
    
    # LLMで要約
    update_dict = {}
    
    # カテゴリごとに処理
    for cat, field_name in [
        ("strategy", "strategy"),
        ("product", "product"),
        ("top", "top"),
        ("industry", "industry"),
        ("society", "society"),
        ("global_items", "global_items")
    ]:
        items = getattr(state, field_name, [])
        
        # 長い文言を要約
        summarized_items = []
        for item in items:
            if len(item) > 20:  # 20文字以上なら要約
                summary = summarize_with_openai(item, cat)
                summarized_items.append(summary)
            else:
                summarized_items.append(item)  # 短い文言はそのまま
        
        update_dict[field_name] = summarized_items
    
    # 更新状態を返す
    result = state.model_copy(update=update_dict)
    return result

def svg_compose(state: SlidesState) -> SlidesState:
    """SVGを生成"""
    # デバッグ出力
    print(f"DEBUG: state.global_items in svg_compose = {state.global_items}")
    
    # テンプレートとパレット読み込み
    template_id = state.template_id
    palette_name = state.palette
    tpl = ENV.get_template(f"{template_id}/slide.svg.j2")
    palette = load_palette(template_id, palette_name)
    
    # コンテンツ設定 - global_itemsをglobalキーとして設定
    content = {
        "company": state.company or "Company",
        "center": state.company or "Company",
        "strategy": state.strategy or [],
        "product": state.product or [],
        "top": state.top or [],
        "industry": state.industry or [],
        "society": state.society or [],
        "global": state.global_items or [],  # 重要: global_itemsをglobalとして渡す
    }
    
    # デバッグ出力
    print(f"DEBUG: グローバルノード用データ = {content['global']}")
    
    # テンプレート変数作成
    ctx = {
        "content": content,
        "palette_colors": palette,
    }
    
    # SVG生成
    svg = tpl.render(**ctx)
    
    # SVGのHTMLエンティティを修正
    svg = fix_svg_entities(svg)
    
    # 更新状態を返す
    return state.model_copy(update={"svg": svg})

def fix_svg_entities(svg: str) -> str:
    """SVGのHTMLエンティティを修正する関数"""
    # エンティティを正しく処理
    svg = re.sub(r'&(?!(amp|lt|gt|quot|apos);)', '&amp;', svg)
    
    # 二重エスケープ防止
    svg = svg.replace('&amp;amp;', '&amp;')
    svg = svg.replace('&amp;lt;', '&lt;')
    svg = svg.replace('&amp;gt;', '&gt;')
    svg = svg.replace('&amp;quot;', '&quot;')
    svg = svg.replace('&amp;apos;', '&apos;')
    
    return svg

# ──────────── Helpers ────────────────
def load_palette(template_id: str, palette_name: str) -> dict:
    """パレットを読み込む"""
    pfile = TPL_DIR / template_id / "palette.json"
    data = json.loads(pfile.read_text(encoding="utf-8"))
    return data[palette_name]

# ──────────── Build graph ─────────────
def create_graph():
    """LangGraphステートマシンを構築"""
    G = StateGraph(SlidesState)
    G.add_node("parse", parse_context)
    G.add_node("summarize", summarize_content)
    G.add_node("compose", svg_compose)
    
    G.add_edge("parse", "summarize")
    G.add_edge("summarize", "compose")
    G.add_edge("compose", END)
    
    G.set_entry_point("parse")
    
    return G.compile()

# ──────────── CLI ─────────────────────
if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", required=True, help="JSON file with company & text")
    ap.add_argument("--export", nargs="*", choices=["pptx"])
    ap.add_argument("--skip-llm", action="store_true", help="LLM処理をスキップ")
    args = ap.parse_args()
    
    # LLMスキップフラグを環境変数に設定
    if args.skip_llm:
        os.environ["SKIP_LLM"] = "true"
    
    # 入力読み込み
    raw = json.loads(Path(args.src).read_text(encoding="utf-8"))
    init = SlidesState(raw_json=raw, company=raw.get("company"))
    
    # グラフ実行
    graph = create_graph()
    result = graph.invoke(init)
    
    # 出力保存
    out = ROOT / "out"; out.mkdir(exist_ok=True)
    svg_path = out / "slide.svg"
    svg_path.write_text(result["svg"], encoding="utf-8")
    print("SVG →", svg_path)
    
    # PPTXエクスポート(オプション)
    if args.export and "pptx" in args.export:
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(svg_path), 0, 0,
                                    prs.slide_width, prs.slide_height)
            ppt_path = out / "slide.pptx"
            prs.save(ppt_path)
            print("PPTX →", ppt_path)
        except ImportError:
            print("python-pptx 未インストールのため PPTX 生成をスキップ")
